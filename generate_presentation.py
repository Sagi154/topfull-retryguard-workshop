"""
Generate the RetryGuard on TopFull project plan presentation.
Run:  python generate_presentation.py
Output: RetryGuard-TopFull-ProjectPlan.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
BLUE   = RGBColor(0x1F, 0x5F, 0x99)
LTBLUE = RGBColor(0xD0, 0xE4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY   = RGBColor(0x55, 0x55, 0x6B)
ACCENT = RGBColor(0xE8, 0x6B, 0x1F)
GREEN  = RGBColor(0x1A, 0x7A, 0x45)
RED    = RGBColor(0xC0, 0x2A, 0x2A)
LTGRAY = RGBColor(0xF4, 0xF6, 0xFA)
MIDGRAY= RGBColor(0xD8, 0xDC, 0xE8)
LGREEN = RGBColor(0xC8, 0xEE, 0xD8)
LRED   = RGBColor(0xF5, 0xD0, 0xD0)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

ARCH_IMG = os.path.join(os.path.dirname(__file__), "Online-Boutique-architecture.png")

# ────────────────────────────────────────────────────────────────────────────
# Core helpers
# ────────────────────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width=Pt(0)):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_width
    else:
        sh.line.fill.background()
    return sh

def add_textbox(slide, left, top, width, height, text,
                font_size=Pt(14), bold=False, color=DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_bullets(slide, left, top, width, height, items,
                font_size=Pt(13), color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    BULLET = "\u2022"
    SUB    = "\u2013"
    for idx, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{SUB if level else BULLET}  {text}"
        run.font.size = font_size if level == 0 else Pt(font_size.pt - 1.5)
        run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.15), NAVY)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)

def add_section_label(slide, label, top=Inches(1.25)):
    add_rect(slide, Inches(0.4), top, Inches(2.2), Inches(0.3), BLUE)
    tb = slide.shapes.add_textbox(Inches(0.45), top + Pt(2), Inches(2.1), Inches(0.28))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9.5)
    run.font.bold = True
    run.font.color.rgb = WHITE

def divider_line(slide, y):
    add_rect(slide, Inches(0.4), y, Inches(12.5), Pt(1), LTBLUE)

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def vert_bar(slide, x):
    """Thin vertical separator between left-text and right-diagram panels."""
    add_rect(slide, x, Inches(1.25), Pt(1.5), Inches(6.0), LTBLUE)

# ────────────────────────────────────────────────────────────────────────────
# Diagram primitives
# ────────────────────────────────────────────────────────────────────────────

def dbox(slide, cx, cy, w, h, text, fill, text_color=WHITE,
         fsize=Pt(11), bold=True, shape_id=5, line_color=None):
    """Centred diagram node at (cx, cy)."""
    sh = slide.shapes.add_shape(shape_id, cx - w/2, cy - h/2, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    lc = line_color if line_color else fill
    sh.line.color.rgb = lc
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top   = Inches(0.03)
    tf.margin_bottom= Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = fsize
    run.font.bold = bold
    run.font.color.rgb = text_color
    return sh

def dline(slide, x1, y1, x2, y2, color=NAVY, width=Pt(2)):
    """Straight connector line."""
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn

def darrow_down(slide, cx, y1, y2, color=NAVY, width=Pt(2)):
    """Vertical arrow downward at horizontal centre cx."""
    dline(slide, cx, y1, cx, y2, color, width)
    # Arrowhead as small triangle (3 thin lines)
    tip = y2
    aw = Inches(0.08)
    dline(slide, cx, tip, cx - aw, tip - aw, color, Pt(1.5))
    dline(slide, cx, tip, cx + aw, tip - aw, color, Pt(1.5))

def darrow_right(slide, x1, x2, cy, color=NAVY, width=Pt(2)):
    """Horizontal arrow rightward at vertical centre cy."""
    dline(slide, x1, cy, x2, cy, color, width)
    tip = x2
    aw = Inches(0.08)
    dline(slide, tip, cy, tip - aw, cy - aw, color, Pt(1.5))
    dline(slide, tip, cy, tip - aw, cy + aw, color, Pt(1.5))

def curved_arrow_left(slide, x1, y1, x2, y2, color=NAVY, width=Pt(2)):
    """Elbow connector (type 2 = elbow) from (x1,y1) to (x2,y2)."""
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = width
    return conn

def dpanel_bg(slide, left, top, width, height, fill=LTBLUE):
    """Lightly shaded background panel for a diagram area."""
    sh = add_rect(slide, left, top, width, height, fill,
                  line_color=MIDGRAY, line_width=Pt(1))
    return sh

def panel_title(slide, left, top, width, text):
    add_textbox(slide, left, top, width, Inches(0.3), text,
                font_size=Pt(9.5), bold=True, color=BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
add_rect(s, 0, 0, W, Inches(4.5), NAVY)
add_rect(s, 0, Inches(4.5), W, Pt(5), ACCENT)

tb = s.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(11.9), Inches(1.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "RetryGuard on TopFull"
run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = WHITE

add_textbox(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(0.75),
    "A Project Plan: Evaluating Dynamic Retry Control on a Kubernetes Microservice Stack",
    font_size=Pt(21), color=RGBColor(0xB0, 0xC8, 0xE8))

add_textbox(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.6),
    "TAU Workshop in Communication Networks and Information Security",
    font_size=Pt(16), color=RGBColor(0x90, 0xB8, 0xD8))

# Mini component diagram on dark background
for i, (label, col) in enumerate([
    ("Locust", BLUE), ("TopFull", NAVY), ("RetryGuard", ACCENT),
    ("Istio / Kubernetes", BLUE), ("Online Boutique", NAVY)]):
    dbox(s, Inches(1.5 + i * 2.25), Inches(5.35), Inches(1.9), Inches(0.52),
         label, col, WHITE, Pt(11))
    if i < 4:
        dline(s, Inches(1.5 + i * 2.25 + 0.95), Inches(5.35),
                 Inches(1.5 + i * 2.25 + 1.3), Inches(5.35), MIDGRAY, Pt(1.5))

add_textbox(s, Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.45),
    "Self-implementation of Algorithm 1 (TAU Deepness Lab, arXiv:2511.23278, 2025)  ·  TopFull SIGCOMM 2024  ·  GCP / Kubernetes 1.26 / Istio",
    font_size=Pt(12), color=GRAY)

note(s, "Cover slide. Project plan presentation.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Goal and Hypothesis
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Goal and Hypothesis", "§1 — What this project is about")
add_section_label(s, "SECTION 1 OF 9")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.3),
    "GOAL", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.7),
    "Self-implement RetryGuard (Algorithm 1, TAU Deepness Lab 2025) and measure its impact when layered on top of TopFull overload control running Online Boutique under overload. We want to know where it helps, where it doesn't, and what trade-offs emerge.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.8))

add_textbox(s, Inches(0.4), Inches(2.95), Inches(12.5), Inches(0.3),
    "HYPOTHESIS", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.25), Inches(12.5), Inches(0.7),
    "Dynamic retry control reduces retry storms and protects goodput compared to default retries. System-wide gain may be modest while specific microservices show large improvement — the experiment is designed to surface both.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(4.1))

add_textbox(s, Inches(0.4), Inches(4.25), Inches(12.5), Inches(0.3),
    "DELIVERABLES", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.0),
    ["Working Kubernetes + Istio + TopFull + RetryGuard experimental setup.",
     "Baseline vs. RetryGuard experiment data across all load scenarios.",
     "Evaluation report with time-series charts comparing goodput, latency, retries per request, resource usage, and autoscaler behavior."],
    font_size=Pt(13))

note(s, "Answers the professor's requirement: 'I want to understand what you are going to do and what the deliverables will be.'")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is TopFull? (Problem + starvation diagram)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "What is TopFull?", "§2 — Problem it solves  ·  SIGCOMM 2024, KAIST")
add_section_label(s, "SECTION 2 OF 9 (1/2)")

# ── Left text column ────────────────────────────────────────────────────────
LW = Inches(6.0)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "THE PROBLEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(1.8),
    "Existing overload controls (DAGOR, Breakwater) manage individual microservices in isolation. When multiple APIs share an overloaded downstream microservice, this causes starvation.\n\nSome API requests are partially processed then rejected downstream, wasting all resources spent on the upstream path.",
    font_size=Pt(12.5), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(3.9), LW, Inches(0.3),
    "HOW TOPFULL SOLVES IT", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.25), LW, Inches(2.6),
    ["Top-down, API-wise control at the entry point — not per-microservice.",
     "Clusters interdependent APIs and solves each cluster independently.",
     "RL-based rate controller (PPO, Sim2Real) — decisions every 1 second.",
     "Respects business priority when throttling and recovering."],
    font_size=Pt(12))

vert_bar(s, Inches(6.65))

# ── Right diagram: starvation node diagram ───────────────────────────────────
RX = Inches(6.9)   # right panel left edge
RC = Inches(9.9)   # right panel centre x

dpanel_bg(s, RX - Inches(0.05), Inches(1.5), Inches(6.35), Inches(5.6))
panel_title(s, RX, Inches(1.55), Inches(6.2), "  STARVATION: THE CORE PROBLEM")

# TopFull box at entry
dbox(s, RC, Inches(2.15), Inches(5.0), Inches(0.52),
     "TopFull Entry Proxy  (throttles API-wise)", NAVY, WHITE, Pt(11))

# API 1 and API 2
dbox(s, RC - Inches(1.7), Inches(3.1), Inches(1.8), Inches(0.45),
     "API 1", BLUE, WHITE, Pt(11))
dbox(s, RC + Inches(1.7), Inches(3.1), Inches(1.8), Inches(0.45),
     "API 2", BLUE, WHITE, Pt(11))

# Arrows: TopFull → API 1 & API 2
darrow_down(s, RC - Inches(1.7), Inches(2.41), Inches(2.87), MIDGRAY, Pt(1.5))
darrow_down(s, RC + Inches(1.7), Inches(2.41), Inches(2.87), MIDGRAY, Pt(1.5))

# Microservice A
dbox(s, RC, Inches(3.85), Inches(5.0), Inches(0.52),
     "Microservice A  (capacity: 10k rps)", NAVY, WHITE, Pt(11))

# Arrows: API 1 & API 2 → MA
darrow_down(s, RC - Inches(1.7), Inches(3.32), Inches(3.59), NAVY, Pt(2))
darrow_down(s, RC + Inches(1.7), Inches(3.32), Inches(3.59), NAVY, Pt(2))

# Microservice B (only API 1's call chain continues)
dbox(s, RC - Inches(1.7), Inches(4.7), Inches(2.6), Inches(0.52),
     "Microservice B\n(capacity: 3k rps)", ACCENT, WHITE, Pt(10))

# Arrow: MA → MB (API 1 path only)
darrow_down(s, RC - Inches(1.7), Inches(4.11), Inches(4.44), ACCENT, Pt(2))

# Label: API 1 path only
add_textbox(s, RC - Inches(0.8), Inches(4.3), Inches(2.2), Inches(0.32),
    "API 1 path only →", font_size=Pt(9), color=ACCENT, bold=False)

# Problem annotation box
add_rect(s, RX, Inches(5.45), Inches(6.25), Inches(1.4),
         RGBColor(0xFF, 0xF0, 0xE0), line_color=ACCENT, line_width=Pt(1))
add_textbox(s, RX + Inches(0.1), Inches(5.5), Inches(6.05), Inches(1.3),
    "Problem: MA wastes capacity on API 1 requests\nthat will be rejected at MB anyway.\n→ API 2 gets less than its fair share of MA: starvation.\nTopFull fixes this by clustering APIs and controlling\nadmission API-wise from the top.",
    font_size=Pt(10.5), color=DARK, wrap=True)

note(s, "The starvation problem is the core motivation for TopFull. The diagram shows why per-microservice controls fail when APIs share downstream services.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is TopFull? (Key Results)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "What is TopFull? — Validation", "§2 — Key results (supporting context)  ·  SIGCOMM 2024")
add_section_label(s, "SECTION 2 OF 9 (2/2)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.45),
    "The following results validate that TopFull works and is worth building on. They establish the baseline system we run RetryGuard on top of.",
    font_size=Pt(13), color=GRAY, italic=True)

divider_line(s, Inches(2.25))

add_textbox(s, Inches(0.4), Inches(2.4), Inches(12.5), Inches(0.3),
    "SIGCOMM 2024 RESULTS (compact, supporting context only)", font_size=Pt(11), bold=True, color=BLUE)

add_bullets(s, Inches(0.4), Inches(2.75), Inches(7.2), Inches(3.5),
    ["1.82× more goodput than DAGOR, 2.26× more than Breakwater during overload.",
     "With autoscaler: up to 3.91× more goodput under traffic surge vs. standalone autoscaler.",
     "Tolerates traffic spikes with up to 57% fewer infrastructure resources.",
     "Converges to optimal rate in 5 seconds vs. 27 seconds for DAGOR."],
    font_size=Pt(13.5))

# ── Right: comparison bar chart (illustrative) ──────────────────────────────
vert_bar(s, Inches(7.6))
dpanel_bg(s, Inches(7.8), Inches(2.1), Inches(5.1), Inches(4.0))
panel_title(s, Inches(7.9), Inches(2.15), Inches(4.9), "  GOODPUT DURING OVERLOAD (illustrative, from TopFull.pdf)")

systems = [("DAGOR",     1.0, MIDGRAY),
           ("Breakwater",0.81, MIDGRAY),
           ("TopFull",   1.82, BLUE)]
bar_h    = Inches(0.45)
bar_maxw = Inches(3.8)
bar_y    = Inches(2.7)
bar_gap  = Inches(0.85)
for label, ratio, col in systems:
    bw = bar_maxw * ratio
    add_rect(s, Inches(8.1), bar_y, bw, bar_h, col)
    add_textbox(s, Inches(8.0) - Inches(1.05), bar_y + Inches(0.08), Inches(1.0), Inches(0.35),
                label, font_size=Pt(10), color=DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, Inches(8.1) + bw + Inches(0.07), bar_y + Inches(0.08), Inches(0.8), Inches(0.35),
                f"{ratio:.2f}×", font_size=Pt(10), bold=True, color=col)
    bar_y += bar_gap

add_textbox(s, Inches(7.9), Inches(5.35), Inches(4.9), Inches(0.5),
    "Normalised to DAGOR goodput = 1.0\nSource: TopFull.pdf (Fig. 8)",
    font_size=Pt(9), color=GRAY, italic=True)

add_textbox(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.35),
    "Source: TopFull.pdf, PRESENTATION-GUIDE.md §2", font_size=Pt(10), color=GRAY, italic=True)

note(s, "Numbers validate TopFull. The experiment question is: does RetryGuard add further value on top of this already-effective system?")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What is RetryGuard? (Problem + retry-storm diagram)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "What is RetryGuard?", "§3 — Problem it solves  ·  TAU Deepness Lab, arXiv:2511.23278, 2025")
add_section_label(s, "SECTION 3 OF 9 (1/2)")

# ── Left text ────────────────────────────────────────────────────────────────
LW = Inches(6.0)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "THE PROBLEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(1.8),
    "Default retry mechanisms (exponential backoff, jitter, retry budgets) are designed for instantaneous failures.\n\nDuring prolonged miscoordination — when services scale at different rates — retries become counterproductive. Each failed retry multiplies load on an already-overloaded service.",
    font_size=Pt(12.5), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(3.85), LW, Inches(0.3),
    "THE CONSEQUENCE", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.2), LW, Inches(2.4),
    ["Rejection rate rises sharply when ρ > 1 (shown analytically, paper Sec. 5).",
     "Each retry multiplies load further — a snowball effect.",
     "Result: self-inflicted Denial-of-Wallet (DoW) — inflated costs, over-scaling, degraded performance.",
     "Effect persists as long as the miscoordination lasts — brief spikes do not trigger it."],
    font_size=Pt(12))

vert_bar(s, Inches(6.65))

# ── Right: retry storm cascade diagram ───────────────────────────────────────
RX = Inches(6.9)
RC = Inches(9.9)

dpanel_bg(s, RX - Inches(0.05), Inches(1.5), Inches(6.35), Inches(5.6))
panel_title(s, RX, Inches(1.55), Inches(6.2), "  THE RETRY STORM FEEDBACK LOOP")

# Upstream service
dbox(s, RC, Inches(2.2), Inches(3.5), Inches(0.52),
     "Upstream Service", BLUE, WHITE, Pt(11))

# Overloaded downstream
dbox(s, RC, Inches(3.55), Inches(3.5), Inches(0.52),
     "Overloaded Service  (ρ > 1)", ACCENT, WHITE, Pt(11))

# Main request arrow
darrow_down(s, RC, Inches(2.46), Inches(3.29), NAVY, Pt(2))
add_textbox(s, RC + Inches(0.22), Inches(2.65), Inches(1.5), Inches(0.3),
    "request", font_size=Pt(10), color=NAVY)

# Rejection arrow back up
dline(s, RC + Inches(1.4), Inches(3.55), RC + Inches(1.4), Inches(2.72),
      RED, Pt(2))
dline(s, RC + Inches(1.4), Inches(2.72), RC + Inches(0.88), Inches(2.72),
      RED, Pt(2))
add_textbox(s, RC + Inches(1.5), Inches(3.0), Inches(1.3), Inches(0.3),
    "rejected!", font_size=Pt(10), bold=True, color=RED)

# Retry arrows (3, wider each time to show amplification)
retry_colors = [RGBColor(0xE8, 0x9B, 0x4F),
                RGBColor(0xE8, 0x6B, 0x1F),
                RGBColor(0xC0, 0x2A, 0x2A)]
retry_offsets = [-Inches(0.55), Inches(0.0), Inches(0.55)]
for i, (rc_off, rcol) in enumerate(zip(retry_offsets, retry_colors)):
    rx1 = RC - Inches(1.2) + rc_off
    ry1 = Inches(2.46)
    ry2 = Inches(3.29)
    dline(s, rx1, ry1, rx1, ry2, rcol, Pt(1.5 + i * 0.8))
    add_textbox(s, rx1 - Inches(0.7), ry1 - Inches(0.28), Inches(0.65), Inches(0.3),
        f"retry {i+1}", font_size=Pt(9), color=rcol)

# Load amplification label
add_rect(s, RX, Inches(4.35), Inches(6.25), Inches(0.45),
         RGBColor(0xFF, 0xF0, 0xE0), line_color=ACCENT, line_width=Pt(1))
add_textbox(s, RX + Inches(0.1), Inches(4.38), Inches(6.05), Inches(0.38),
    "Each retry adds to ρ  →  more rejections  →  more retries  →  ρ rises further",
    font_size=Pt(11), bold=True, color=ACCENT)

# DoW result
dbox(s, RC, Inches(5.2), Inches(5.5), Inches(0.62),
     "Denial-of-Wallet: over-scaling, inflated billing, degraded SLO",
     RED, WHITE, Pt(10.5))

# Arrow from overloaded → DoW
darrow_down(s, RC, Inches(4.07), Inches(4.89), RED, Pt(2))

note(s, "The retry storm is a positive feedback loop. The key insight of RetryGuard is that prolonged overload (ρ > 1 sustained) is analytically distinguishable from transient failures.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What is RetryGuard? (Algorithm 1 state machine)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "What is RetryGuard? — Algorithm 1 & Integration", "§3 — How it works and how we run it")
add_section_label(s, "SECTION 3 OF 9 (2/2)")

# ── Left text ────────────────────────────────────────────────────────────────
LW = Inches(5.8)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "ALGORITHM 1 — PER-SERVICE RETRY CONTROLLER", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(2.0), LW, Inches(2.7),
    ["Monitors rejection rate per service continuously.",
     "If rejection rate > ~20% threshold for Interval consecutive periods (~30s each) → Consecutive_high ≥ Interval: DISABLE retries for that service.",
     "If rejection rate < threshold for Interval consecutive periods → Consecutive_low ≥ Interval: RE-ENABLE retries.",
     "Distributed — each service runs its own instance. No central orchestrator.",
     "Non-intrusive under normal load (ρ < 1) — false positives are negligible (paper Fig. 7)."],
    font_size=Pt(12))

divider_line(s, Inches(4.9))

add_textbox(s, Inches(0.4), Inches(5.05), LW, Inches(0.3),
    "ISTIO INTEGRATION (Appendix A + Sec. 4)", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(5.4), LW, Inches(1.7),
    ["Python script on master node samples Istio/Envoy sidecar HTTP error rates.",
     "Patches Istio VirtualService retry policy per service via Kubernetes Python client.",
     "No per-request overhead — 30s polling cycle. No application code changes."],
    font_size=Pt(12))

vert_bar(s, Inches(6.45))

# ── Right: Algorithm 1 state machine ─────────────────────────────────────────
RX = Inches(6.65)
RC = Inches(9.95)

dpanel_bg(s, RX - Inches(0.05), Inches(1.5), Inches(6.6), Inches(5.6))
panel_title(s, RX, Inches(1.55), Inches(6.4), "  ALGORITHM 1 STATE MACHINE  (per service, independently)")

STATE_Y1 = Inches(2.8)
STATE_Y2 = Inches(5.2)
STATE_W  = Inches(3.0)
STATE_H  = Inches(0.75)

# RETRIES ON state
on_box = dbox(s, RC, STATE_Y1, STATE_W, STATE_H,
              "RETRIES ON", GREEN, WHITE, Pt(16), shape_id=5)

# RETRIES OFF state
off_box = dbox(s, RC, STATE_Y2, STATE_W, STATE_H,
               "RETRIES OFF", RED, WHITE, Pt(16), shape_id=5)

# Monitoring cycle annotation inside each state
add_textbox(s, RC - STATE_W/2 + Inches(0.1), STATE_Y1 + Inches(0.02),
            STATE_W - Inches(0.2), Inches(0.28),
            "normal operation — no intervention", font_size=Pt(8.5),
            color=LGREEN, align=PP_ALIGN.CENTER)

add_textbox(s, RC - STATE_W/2 + Inches(0.1), STATE_Y2 + Inches(0.02),
            STATE_W - Inches(0.2), Inches(0.28),
            "VirtualService CRD: retries disabled", font_size=Pt(8.5),
            color=LRED, align=PP_ALIGN.CENTER)

# ── Downward transition (ON → OFF): right side ──────────────────────────────
TX_R = RC + Inches(2.0)   # x for right-side elbow

dline(s, RC + STATE_W/2, STATE_Y1 + STATE_H/2,
         TX_R, STATE_Y1 + STATE_H/2, ACCENT, Pt(2))
dline(s, TX_R, STATE_Y1 + STATE_H/2,
         TX_R, STATE_Y2 + STATE_H/2, ACCENT, Pt(2))
dline(s, TX_R, STATE_Y2 + STATE_H/2,
         RC + STATE_W/2, STATE_Y2 + STATE_H/2, ACCENT, Pt(2))
# arrowhead hint
dline(s, RC + STATE_W/2, STATE_Y2 + STATE_H/2,
         RC + STATE_W/2 + Inches(0.1), STATE_Y2 + STATE_H/2 - Inches(0.08), ACCENT, Pt(1.5))
dline(s, RC + STATE_W/2, STATE_Y2 + STATE_H/2,
         RC + STATE_W/2 + Inches(0.1), STATE_Y2 + STATE_H/2 + Inches(0.08), ACCENT, Pt(1.5))

add_textbox(s, TX_R + Inches(0.08), (STATE_Y1 + STATE_Y2) / 2 - Inches(0.45),
            Inches(1.5), Inches(0.9),
            "Consecutive_high\n≥ Interval\n(rejection > 20%\nfor N × 30s)",
            font_size=Pt(9), color=ACCENT, wrap=True)

# ── Upward transition (OFF → ON): left side ──────────────────────────────────
TX_L = RC - Inches(2.0)

dline(s, RC - STATE_W/2, STATE_Y2 + STATE_H/2,
         TX_L, STATE_Y2 + STATE_H/2, GREEN, Pt(2))
dline(s, TX_L, STATE_Y2 + STATE_H/2,
         TX_L, STATE_Y1 + STATE_H/2, GREEN, Pt(2))
dline(s, TX_L, STATE_Y1 + STATE_H/2,
         RC - STATE_W/2, STATE_Y1 + STATE_H/2, GREEN, Pt(2))
dline(s, RC - STATE_W/2, STATE_Y1 + STATE_H/2,
         RC - STATE_W/2 - Inches(0.1), STATE_Y1 + STATE_H/2 - Inches(0.08), GREEN, Pt(1.5))
dline(s, RC - STATE_W/2, STATE_Y1 + STATE_H/2,
         RC - STATE_W/2 - Inches(0.1), STATE_Y1 + STATE_H/2 + Inches(0.08), GREEN, Pt(1.5))

add_textbox(s, TX_L - Inches(1.55), (STATE_Y1 + STATE_Y2) / 2 - Inches(0.45),
            Inches(1.5), Inches(0.9),
            "Consecutive_low\n≥ Interval\n(rejection < 20%\nfor N × 30s)",
            font_size=Pt(9), color=GREEN, wrap=True)

# ── Monitoring loop label ────────────────────────────────────────────────────
add_textbox(s, RC - Inches(2.6), Inches(3.95), Inches(5.2), Inches(0.38),
    "Sample Istio/Envoy sidecar metrics  ·  ~every 30 s", 
    font_size=Pt(9.5), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# Key results footnote
add_rect(s, RX, Inches(6.35), Inches(6.6), Inches(0.75),
         LTBLUE, line_color=MIDGRAY, line_width=Pt(1))
add_textbox(s, RX + Inches(0.1), Inches(6.4), Inches(6.4), Inches(0.65),
    "Results (TAU 2025): AWS retries 2.09→0.05/req (−98%), billing 1029%→100%  ·  Istio/K8s retries 0.31→0.01/req, billing 224%→100%",
    font_size=Pt(9.5), color=DARK)

note(s, "The state machine is exactly Algorithm 1. Each service has its own independent instance. The Consecutive_high and Consecutive_low counters prevent oscillation.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Stack & Topology
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Stack & Topology", "§4 — Infrastructure and component layout")
add_section_label(s, "SECTION 4 OF 9")

# ── Left: request-path flow diagram ──────────────────────────────────────────
dpanel_bg(s, Inches(0.3), Inches(1.5), Inches(6.3), Inches(5.7))
panel_title(s, Inches(0.4), Inches(1.55), Inches(6.1), "  REQUEST PATH & CONTROL LOOPS")

flow_cx = Inches(3.4)
components = [
    ("Locust  (load-gen VM)",          GRAY,  Inches(2.0)),
    ("TopFull Rate Limiter\n(master)",  NAVY,  Inches(2.65)),
    ("Istio / Envoy sidecars",          BLUE,  Inches(3.3)),
    ("Online Boutique\nmicroservices",  NAVY,  Inches(4.05)),
]
comp_y = [Inches(2.1), Inches(3.0), Inches(3.95), Inches(4.9)]
for (label, col, _), cy in zip(components, comp_y):
    dbox(s, flow_cx, cy, Inches(3.5), Inches(0.55), label, col, WHITE, Pt(10.5))

for i in range(len(comp_y) - 1):
    darrow_down(s, flow_cx, comp_y[i] + Inches(0.27),
                comp_y[i+1] - Inches(0.27), MIDGRAY, Pt(2))

# RetryGuard side box
dbox(s, Inches(1.3), Inches(3.95), Inches(1.7), Inches(0.7),
     "Retry-\nGuard", ACCENT, WHITE, Pt(10.5))
dline(s, Inches(2.15), Inches(3.95), Inches(1.85) + Inches(0.55), Inches(3.95),
      ACCENT, Pt(1.5))
add_textbox(s, Inches(0.35), Inches(4.65), Inches(2.5), Inches(0.35),
    "patches VirtualService\nCRDs (Kubernetes API)", font_size=Pt(8.5), color=ACCENT)

# RL feedback arrow (TopFull)
dline(s, Inches(5.15), Inches(3.0), Inches(5.55), Inches(3.0), BLUE, Pt(1.5))
dline(s, Inches(5.55), Inches(3.0), Inches(5.55), Inches(4.9), BLUE, Pt(1.5))
dline(s, Inches(5.55), Inches(4.9), Inches(5.15), Inches(4.9), BLUE, Pt(1.5))
add_textbox(s, Inches(5.6), Inches(3.7), Inches(1.0), Inches(0.8),
    "RL\nfeed-\nback\n(1 s)", font_size=Pt(8.5), color=BLUE)

# VM roles footnote
add_rect(s, Inches(0.35), Inches(6.35), Inches(6.2), Inches(0.65),
         LTBLUE, line_color=MIDGRAY, line_width=Pt(1))
add_textbox(s, Inches(0.45), Inches(6.4), Inches(6.0), Inches(0.55),
    "Master: 8+ vCPU, 16GB — K8s ctrl, Istio, TopFull, RetryGuard  ·  Workers: pods, cAdvisor  ·  Load-gen: Locust only  ·  GCP (~$300 credits)",
    font_size=Pt(9.5), color=DARK)

# ── Right: Online Boutique architecture image ────────────────────────────────
if os.path.exists(ARCH_IMG):
    s.shapes.add_picture(ARCH_IMG, Inches(6.8), Inches(1.35), Inches(6.1), Inches(5.5))
add_textbox(s, Inches(6.8), Inches(6.9), Inches(6.1), Inches(0.35),
    "Online Boutique call graph — 11 microservices. Representative test app, not the subject of study.",
    font_size=Pt(9.5), color=GRAY, italic=True)

note(s, "RetryGuard on master: reads Istio sidecar metrics (separate from TopFull's entry-proxy collectors), patches VirtualService CRDs.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — How We Test: The Baseline
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "How We Test — The Baseline", "§5 — Controlled experimental methodology (1/3)")
add_section_label(s, "SECTION 5 OF 9 (1/3)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.55),
    "To isolate RetryGuard's effect, we need a baseline that captures the system's behavior without it. Everything except the retry policy must be identical between the baseline run and the RetryGuard run.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.35))

add_textbox(s, Inches(0.4), Inches(2.5), Inches(7.2), Inches(0.3),
    "BASELINE CONFIGURATION", font_size=Pt(11), bold=True, color=BLUE)

add_bullets(s, Inches(0.4), Inches(2.85), Inches(7.2), Inches(3.8),
    ["TopFull overload controller: ON.",
     "Istio default retry policy: ON. Retries fire per default VirtualService config.",
     "RetryGuard: OFF. Controller script not running.",
     "Fixed workload scenario — same load shape, duration, replica counts every run.",
     "Output → named baseline artifact dir: CSVs from metric_collector.py, resource_collector.py, overload_detection.py.",
     "Repeated multiple times to account for Locust non-determinism."],
    font_size=Pt(13))

# ── Right: config summary card ────────────────────────────────────────────────
dpanel_bg(s, Inches(7.8), Inches(2.35), Inches(5.1), Inches(4.7))
panel_title(s, Inches(7.9), Inches(2.4), Inches(4.9), "  BASELINE — COMPONENT STATES")

components_bl = [
    ("TopFull",           "ON",  GREEN),
    ("Locust workload",   "ON",  GREEN),
    ("Istio retries",     "ON (default)", GREEN),
    ("RetryGuard",        "OFF", RED),
]
cy = Inches(3.0)
for comp, state, col in components_bl:
    dbox(s, Inches(9.35), cy, Inches(3.6), Inches(0.52), comp, NAVY, WHITE, Pt(12))
    dbox(s, Inches(11.55), cy, Inches(1.1), Inches(0.52), state, col, WHITE, Pt(10.5))
    cy += Inches(0.8)

add_textbox(s, Inches(7.9), Inches(6.0), Inches(4.9), Inches(0.55),
    "Single variable changes between baseline\nand experiment: RetryGuard ON vs. OFF.",
    font_size=Pt(10.5), color=NAVY, bold=True)

note(s, "The baseline is the control condition: TopFull doing its best without RetryGuard. All improvements are measured as deltas from this.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — How We Test: The Experiment
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "How We Test — The Experiment", "§5 — Controlled experimental methodology (2/3)")
add_section_label(s, "SECTION 5 OF 9 (2/3)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.55),
    "The RetryGuard experiment is identical to the baseline in every respect except one: RetryGuard is running. This single-variable change lets us attribute any outcome difference directly to RetryGuard.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.35))

add_textbox(s, Inches(0.4), Inches(2.5), Inches(7.2), Inches(0.3),
    "EXPERIMENT CONFIGURATION", font_size=Pt(11), bold=True, color=BLUE)

add_bullets(s, Inches(0.4), Inches(2.85), Inches(7.2), Inches(3.8),
    ["Same load, topology, replica counts as baseline — no other changes.",
     "RetryGuard: ON. Python script on master node, monitoring Istio/Envoy metrics.",
     "Controller parameters from Sec. 6.2: ~20% rejection threshold, ~30s interval.",
     "Patches Istio VirtualService CRDs via Kubernetes Python client when overload detected.",
     "Output → named experiment artifact dir mirroring baseline structure.",
     "Controller decisions logged: service name, toggle direction, time, rejection rate reading."],
    font_size=Pt(13))

# ── Right: config summary card ────────────────────────────────────────────────
dpanel_bg(s, Inches(7.8), Inches(2.35), Inches(5.1), Inches(4.7))
panel_title(s, Inches(7.9), Inches(2.4), Inches(4.9), "  EXPERIMENT — COMPONENT STATES")

components_ex = [
    ("TopFull",           "ON",  GREEN),
    ("Locust workload",   "ON",  GREEN),
    ("Istio retries",     "managed", BLUE),
    ("RetryGuard",        "ON", GREEN),
]
cy = Inches(3.0)
for comp, state, col in components_ex:
    dbox(s, Inches(9.35), cy, Inches(3.6), Inches(0.52), comp, NAVY, WHITE, Pt(12))
    dbox(s, Inches(11.55), cy, Inches(1.1), Inches(0.52), state, col, WHITE, Pt(10.5))
    cy += Inches(0.8)

add_rect(s, Inches(7.9), Inches(5.95), Inches(4.9), Inches(0.7),
         RGBColor(0xFF, 0xF0, 0xE0), line_color=ACCENT, line_width=Pt(1))
add_textbox(s, Inches(8.0), Inches(6.0), Inches(4.7), Inches(0.6),
    "Only change from baseline: RetryGuard ON.\nAll other variables identical.",
    font_size=Pt(11), bold=True, color=ACCENT)

note(s, "One variable changes. This is why we can attribute differences in goodput, latency, and resource usage to RetryGuard's presence.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — How We Test: Repeated Runs
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "How We Test — Repeated Runs", "§5 — Controlled experimental methodology (3/3)")
add_section_label(s, "SECTION 5 OF 9 (3/3)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.55),
    "Locust generates randomized user behavior — the arrival pattern is non-deterministic. A single run per scenario is insufficient to isolate RetryGuard's effect from natural traffic variation.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.35))

add_textbox(s, Inches(0.4), Inches(2.5), Inches(12.5), Inches(0.3),
    "REPEATED-RUN PROTOCOL", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(2.85), Inches(7.5), Inches(2.8),
    ["Both baseline and RetryGuard experiment run multiple times per scenario.",
     "Results compared using averages and medians across runs — not individual trial outcomes.",
     "Same inputs, different retry policy, multiple runs: isolates RetryGuard from Locust noise.",
     "Inconsistent runs are flagged for re-investigation before drawing conclusions."],
    font_size=Pt(13))

divider_line(s, Inches(5.0))

add_textbox(s, Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.3),
    "COMPARISON LOGIC", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.7),
    "For each scenario: [Baseline runs] vs. [RetryGuard runs] → delta in goodput, latency, retries/request, CPU, pod count. The delta is the measured effect of RetryGuard. We report both system-wide aggregates and per-microservice breakdowns.",
    font_size=Pt(13.5), color=DARK, wrap=True)

note(s, "Non-determinism in Locust is expected and accounted for by design.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What We Want to Find Out: Opening Frame (with 2-loop diagram)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "What We Want to Find Out", "§6 — The intellectual core of the project (1/3)")
add_section_label(s, "SECTION 6 OF 9 (1/3)")

# ── Left: text ────────────────────────────────────────────────────────────────
LW = Inches(5.9)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "THE CENTRAL QUESTION", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(1.1),
    "Does adding RetryGuard on top of TopFull actually make things better — or is TopFull's overload control already sufficient on its own?",
    font_size=Pt(16), bold=True, color=NAVY, wrap=True)

add_textbox(s, Inches(0.4), Inches(3.2), LW, Inches(3.4),
    "RetryGuard has been validated in AWS Lambda/DynamoDB and standalone Istio/Kubernetes (Bookinfo). It has never been evaluated alongside a sophisticated top-down overload controller.\n\nTopFull controls admission at the entry proxy. But once a request is admitted, Istio's internal retry policies fire inside the cluster — invisible to TopFull's rate limiter. These internal retries can amplify load on downstream services beyond what TopFull can throttle.\n\nRetryGuard, operating per-service with Istio metrics, is positioned to suppress exactly this internal amplification. Whether that suppression produces measurable gains on top of TopFull — and for which services — is the gap this project investigates.",
    font_size=Pt(12.5), color=DARK, wrap=True)

vert_bar(s, Inches(6.55))

# ── Right: two-loop architecture diagram ──────────────────────────────────────
RX = Inches(6.75)
RW = Inches(6.45)
RC = RX + RW / 2

dpanel_bg(s, RX - Inches(0.05), Inches(1.5), RW + Inches(0.1), Inches(5.7))
panel_title(s, RX, Inches(1.55), RW, "  TWO FEEDBACK LOOPS RUNNING SIMULTANEOUSLY")

# Outer loop label: TopFull (1s)
add_rect(s, RX + Inches(0.05), Inches(2.0), RW - Inches(0.1), Inches(4.85),
         RGBColor(0xEA, 0xF4, 0xFF), line_color=BLUE, line_width=Pt(1.5))
add_textbox(s, RX + Inches(0.2), Inches(2.05), Inches(2.5), Inches(0.28),
    "TopFull loop  (1-second cycle)", font_size=Pt(9.5), bold=True, color=BLUE)

# Inner box: Kubernetes cluster
add_rect(s, RX + Inches(0.55), Inches(2.8), RW - Inches(1.1), Inches(3.4),
         RGBColor(0xF0, 0xF4, 0xFF), line_color=NAVY, line_width=Pt(1.5))
add_textbox(s, RX + Inches(0.7), Inches(2.85), Inches(2.5), Inches(0.28),
    "Kubernetes cluster", font_size=Pt(9.5), bold=True, color=NAVY)

# Entry proxy
dbox(s, RC, Inches(2.38), Inches(3.8), Inches(0.48),
     "TopFull Entry Proxy  (rate limiter)", NAVY, WHITE, Pt(10))

# Admitted requests arrow
darrow_down(s, RC, Inches(2.62), Inches(2.92), NAVY, Pt(2))
add_textbox(s, RC + Inches(0.1), Inches(2.68), Inches(1.8), Inches(0.28),
    "admitted requests", font_size=Pt(8.5), color=GRAY)

# Microservices (3 boxes)
svc_cx = [RC - Inches(1.5), RC, RC + Inches(1.5)]
svc_labels = ["Svc A", "Svc B", "Svc C"]
svc_y = Inches(3.5)
for cx, lbl in zip(svc_cx, svc_labels):
    dbox(s, cx, svc_y, Inches(0.95), Inches(0.48), lbl, BLUE, WHITE, Pt(10))

# Internal retry arrows between services
dline(s, svc_cx[0] + Inches(0.47), svc_y,
         svc_cx[1] - Inches(0.47), svc_y, ACCENT, Pt(1.5))
dline(s, svc_cx[1] + Inches(0.47), svc_y,
         svc_cx[2] - Inches(0.47), svc_y, ACCENT, Pt(1.5))
add_textbox(s, RC - Inches(1.3), svc_y - Inches(0.38), Inches(2.6), Inches(0.3),
    "← internal Istio retries (invisible to TopFull) →",
    font_size=Pt(8.5), color=ACCENT, align=PP_ALIGN.CENTER)

# RetryGuard box inside cluster
dbox(s, RC, Inches(4.7), Inches(3.6), Inches(0.52),
     "RetryGuard  (master node, 30s cycle)", ACCENT, WHITE, Pt(10))

# RetryGuard reads sidecar metrics
dline(s, RC - Inches(0.7), Inches(4.46), RC - Inches(0.7), Inches(4.44), ACCENT, Pt(1.5))
curved_arrow_left(s, RC - Inches(1.3), svc_y + Inches(0.24),
                     RC - Inches(1.3), Inches(4.7), ACCENT, Pt(1.5))
add_textbox(s, RX + Inches(0.1), Inches(4.3), Inches(1.8), Inches(0.32),
    "reads sidecar\nmetrics", font_size=Pt(8.5), color=ACCENT)

# TopFull feedback arrow (right side of outer box)
dline(s, RX + RW - Inches(0.25), Inches(2.38),
         RX + RW - Inches(0.25), Inches(3.98), BLUE, Pt(1.5))
dline(s, RX + RW - Inches(0.25), Inches(3.98),
         RC + Inches(1.9), Inches(3.98), BLUE, Pt(1.5))
add_textbox(s, RX + RW - Inches(0.2), Inches(2.95), Inches(0.9), Inches(0.7),
    "RL feed-\nback\n1 s", font_size=Pt(8.5), color=BLUE)

note(s, "The key insight: both TopFull and RetryGuard are feedback loops on overlapping signals. Whether they cooperate or interfere is one of the core open questions.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Open Questions (1/2)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Open Questions", "§6 — Specific questions the project is designed to answer (2/3)")
add_section_label(s, "SECTION 6 OF 9 (2/3)")

questions_1 = [
    ("System-level gains",
     "Does RetryGuard further improve global goodput and latency during overload, or does TopFull's entry-point control already absorb the retry problem so that internal retries are no longer a meaningful factor?"),
    ("Topology beneficiaries",
     "Which specific microservices in the call chain benefit most? Do leaf-node services respond differently than gateway-adjacent services? Is the benefit uneven across the topology?"),
    ("Chain propagation",
     "If RetryGuard activates on one downstream service, do the resource savings propagate upward through the rest of the execution path — or is the effect local to that service?"),
    ("Controller interaction",
     "TopFull's RL controller and RetryGuard are both feedback loops running simultaneously on overlapping signals. TopFull adjusts admission rates every 1 second; RetryGuard toggles per-service retry policies on a ~30-second cycle. Do they cooperate — or does one loop's correction interfere with the other's? This combination has not been studied."),
]

y = Inches(1.65)
for title, body in questions_1:
    add_textbox(s, Inches(0.4), y, Inches(12.5), Inches(0.28),
                title.upper(), font_size=Pt(10.5), bold=True, color=ACCENT)
    add_textbox(s, Inches(0.4), y + Inches(0.3), Inches(12.5), Inches(0.6),
                body, font_size=Pt(12.5), color=DARK, wrap=True)
    y += Inches(1.05)

note(s, "Each question maps to one or more load scenarios. The questions motivate the scenarios — not the other way around.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Open Questions (2/2)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Open Questions (continued)", "§6 — Specific questions the project is designed to answer (3/3)")
add_section_label(s, "SECTION 6 OF 9 (3/3)")

questions_2 = [
    ("Combined equilibrium",
     "When RetryGuard suppresses retries at a bottleneck, the bottleneck's load drops, improving goodput and latency signals that TopFull's RL observes. TopFull may respond by increasing admission. Does this find a better stable throughput point — or does increased admission re-trigger overload and undo the gains?"),
    ("Topology position sensitivity",
     "Does the structural position of the bottleneck service change RetryGuard's contribution? Three positions: (1) gateway-adjacent (e.g., ProductCatalog) — TopFull sees this most directly; (2) hub (e.g., Checkout) — fans out to 6 downstream callers, creating bidirectional retry amplification; (3) deep leaf (e.g., Email or Payment) — most attenuated from TopFull's entry signal. The hub case is expected to show the most severe amplification."),
    ("Interval parameter sensitivity",
     "Is RetryGuard's 30-second re-enable interval optimal when TopFull is co-running? That interval was validated without a simultaneous top-down controller. With TopFull's RL making decisions every 1 second, recovery dynamics may be faster or more oscillatory. Does the optimal interval shift?"),
    ("Adversarial resilience",
     "Under malicious traffic, does hostile load trick the controller into misfiring — suppressing retries when it shouldn't — or does RetryGuard successfully blunt retry amplification from attack traffic?"),
]

y = Inches(1.65)
for title, body in questions_2:
    add_textbox(s, Inches(0.4), y, Inches(12.5), Inches(0.28),
                title.upper(), font_size=Pt(10.5), bold=True, color=ACCENT)
    add_textbox(s, Inches(0.4), y + Inches(0.3), Inches(12.5), Inches(0.65),
                body, font_size=Pt(12), color=DARK, wrap=True)
    y += Inches(1.15)

note(s, "The interval sensitivity and combined equilibrium questions are unique to the TopFull+RetryGuard combination.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Load Scenarios Intro
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Load Scenarios — Overview", "§7 — Each scenario operationalizes one or more open questions")
add_section_label(s, "SECTION 7 OF 9 (INTRO)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.65),
    "The load scenarios are derived directly from the open questions above — they were not chosen independently. Each scenario is designed to answer specific questions. All scenarios use TopFull's built-in synthetic workload generator (Locust + TopFull scripts).",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.45))

# ── Scenario → Question mapping visual ──────────────────────────────────────
add_textbox(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(0.3),
    "SCENARIO → OPEN QUESTION MAPPING", font_size=Pt(11), bold=True, color=BLUE)

rows = [
    ("Normal Operation",          "System-level gains (sanity check)"),
    ("Sustained Overload",        "System-level gains, topology beneficiaries, chain propagation, controller interaction"),
    ("Targeted Bottleneck",       "Topology beneficiaries, chain propagation, controller interaction"),
    ("Topology Position Compare", "Topology position sensitivity, topology beneficiaries, chain propagation"),
    ("Re-enable Interval Tuning", "Interval parameter sensitivity, combined equilibrium"),
    ("Attack Traffic (extension)","Adversarial resilience"),
]
cy = Inches(2.95)
cols = [NAVY, BLUE, BLUE, NAVY, BLUE, GRAY]
for (scen, answers), col in zip(rows, cols):
    add_rect(s, Inches(0.4), cy, Inches(3.7), Inches(0.4), col)
    add_textbox(s, Inches(0.5), cy + Inches(0.05), Inches(3.5), Inches(0.35),
                scen, font_size=Pt(11), bold=True, color=WHITE)
    dline(s, Inches(4.1), cy + Inches(0.2), Inches(4.5), cy + Inches(0.2), MIDGRAY, Pt(1.5))
    dline(s, Inches(4.5), cy + Inches(0.2), Inches(4.4), cy + Inches(0.12), MIDGRAY, Pt(1.5))
    dline(s, Inches(4.5), cy + Inches(0.2), Inches(4.4), cy + Inches(0.28), MIDGRAY, Pt(1.5))
    add_textbox(s, Inches(4.6), cy + Inches(0.05), Inches(8.1), Inches(0.35),
                answers, font_size=Pt(11), color=DARK)
    cy += Inches(0.57)

add_textbox(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
    "RetryGuard per-service decisions are driven by Istio/Envoy sidecar metrics — a different measurement point from TopFull's entry-proxy collectors. Both data streams are collected and cross-referenced.",
    font_size=Pt(11.5), color=GRAY, italic=True)

note(s, "The scenarios follow from the questions — not the other way around. This is the key structural argument for scientific credibility.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Scenario: Normal Operation
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Normal Operation", "§7 — Load scenario 1 of 6")
add_section_label(s, "SECTION 7 OF 9 (1/6)")

LW = Inches(7.5)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "WHAT HAPPENS IN THE SYSTEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(0.7),
    "Traffic is flat and manageable — RPS well within service capacity. No service reaches ρ > 1. TopFull's rate limiter is not under pressure. Istio retry policies fire occasionally for transient errors but do not accumulate.",
    font_size=Pt(13), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(2.8), LW, Inches(0.3),
    "WHAT IT TESTS", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.1), LW, Inches(0.75),
    "Does RetryGuard stay entirely non-intrusive when the system is healthy? The controller should detect rejection rates below the ~20% threshold across all services and leave Istio VirtualService configurations untouched. Spurious retry suppression under healthy load = controller malfunction.",
    font_size=Pt(13), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(4.0), LW, Inches(0.3),
    "OPEN QUESTION ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(4.3), LW, Inches(0.45),
    "System-level gains (non-overload side): necessary sanity check — performance must be indistinguishable between baseline and RetryGuard runs during normal load.",
    font_size=Pt(13), color=DARK, wrap=True)

vert_bar(s, Inches(8.1))

# ── Right: load profile sketch ───────────────────────────────────────────────
dpanel_bg(s, Inches(8.3), Inches(1.5), Inches(4.6), Inches(3.4))
panel_title(s, Inches(8.4), Inches(1.55), Inches(4.4), "  LOAD PROFILE")

# Axes
dline(s, Inches(8.65), Inches(4.5), Inches(8.65), Inches(2.1), GRAY, Pt(1.5))
dline(s, Inches(8.65), Inches(4.5), Inches(12.65), Inches(4.5), GRAY, Pt(1.5))
add_textbox(s, Inches(8.3), Inches(2.05), Inches(0.35), Inches(0.3),
    "RPS", font_size=Pt(8), color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(12.5), Inches(4.55), Inches(0.6), Inches(0.25),
    "time", font_size=Pt(8), color=GRAY)

# Capacity line (ρ = 1)
dline(s, Inches(8.65), Inches(3.15), Inches(12.6), Inches(3.15),
      RGBColor(0xCC, 0x44, 0x44), Pt(1.5))
add_textbox(s, Inches(8.68), Inches(2.9), Inches(1.0), Inches(0.25),
    "ρ = 1", font_size=Pt(8.5), bold=True, color=RED)

# Flat load line (well below capacity)
dline(s, Inches(8.65), Inches(4.0), Inches(12.6), Inches(4.0),
      GREEN, Pt(3))
add_textbox(s, Inches(10.0), Inches(3.65), Inches(2.0), Inches(0.3),
    "Normal load", font_size=Pt(9.5), color=GREEN, bold=True)

# RetryGuard stays off label
add_rect(s, Inches(8.35), Inches(4.65), Inches(4.5), Inches(0.35), LGREEN,
         line_color=GREEN, line_width=Pt(1))
add_textbox(s, Inches(8.45), Inches(4.68), Inches(4.3), Inches(0.3),
    "RetryGuard: no intervention (rejection rate < 20%)", font_size=Pt(9), color=GREEN)

note(s, "Validates non-intrusiveness. The paper claims ρ < 1 is analytically distinguishable — this is the empirical check in our setup.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Scenario: Sustained Overload (with load profile diagram)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Sustained Overload", "§7 — Load scenario 2 of 6  ·  The core experiment")
add_section_label(s, "SECTION 7 OF 9 (2/6)")

LW = Inches(6.8)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "WHAT HAPPENS IN THE SYSTEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(1.05),
    "Load increases until ρ > 1 system-wide and holds there for several minutes. TopFull throttles at the entry proxy. But once a request is admitted, Istio's default retry policy fires internally on each rejection — invisible to TopFull's rate limiter, amplifying load on downstream services.",
    font_size=Pt(12.5), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(3.15), LW, Inches(0.3),
    "WHAT IT TESTS", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.45), LW, Inches(0.7),
    "Does RetryGuard suppress this internal retry amplification — producing measurable improvement in goodput or resource usage on top of what TopFull achieves alone? RetryGuard's ~30s detection window must trigger before any effect is observable.",
    font_size=Pt(12.5), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(4.3), LW, Inches(0.3),
    "OPEN QUESTIONS ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.65), LW, Inches(1.75),
    ["System-level gains — primary outcome.",
     "Topology beneficiaries and chain propagation.",
     "Controller interaction — do TopFull and RetryGuard cooperate or interfere?"],
    font_size=Pt(12.5))

vert_bar(s, Inches(7.35))

# ── Right: load + RetryGuard activation timeline ─────────────────────────────
dpanel_bg(s, Inches(7.55), Inches(1.5), Inches(5.4), Inches(5.6))
panel_title(s, Inches(7.65), Inches(1.55), Inches(5.2),
    "  LOAD PROFILE + RETRYGUARD ACTIVATION")

AX = Inches(7.9)   # axis left
AY = Inches(6.6)   # axis bottom
AW = Inches(4.8)   # axis width
AH = Inches(4.7)   # axis height

# Axes
dline(s, AX, AY, AX, AY - AH, GRAY, Pt(1.5))
dline(s, AX, AY, AX + AW, AY, GRAY, Pt(1.5))
add_textbox(s, AX - Inches(0.45), AY - AH - Inches(0.15), Inches(0.45), Inches(0.3),
    "RPS", font_size=Pt(8), color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, AX + AW + Inches(0.05), AY, Inches(0.5), Inches(0.25),
    "time", font_size=Pt(8), color=GRAY)

# Capacity line
cap_y = AY - AH * 0.55
dline(s, AX, cap_y, AX + AW, cap_y, RED, Pt(1.2))
add_textbox(s, AX - Inches(0.45), cap_y - Inches(0.15), Inches(0.45), Inches(0.25),
    "ρ=1", font_size=Pt(7.5), bold=True, color=RED, align=PP_ALIGN.CENTER)

# Load line: rises and stays above ρ=1
p1 = (AX,             AY - AH * 0.20)   # start: low load
p2 = (AX + AW * 0.25, AY - AH * 0.20)  # flat start
p3 = (AX + AW * 0.40, AY - AH * 0.75)  # ramp up
p4 = (AX + AW,        AY - AH * 0.75)  # sustained overload

dline(s, p1[0], p1[1], p2[0], p2[1], NAVY, Pt(2.5))
dline(s, p2[0], p2[1], p3[0], p3[1], NAVY, Pt(2.5))
dline(s, p3[0], p3[1], p4[0], p4[1], ACCENT, Pt(2.5))

add_textbox(s, p3[0] + Inches(0.1), p3[1] - Inches(0.35), Inches(1.5), Inches(0.3),
    "ρ > 1 sustained", font_size=Pt(8.5), bold=True, color=ACCENT)

# RetryGuard detection window (shaded area after ~30s)
det_x = p3[0] + AW * 0.18
det_end = p4[0]
add_rect(s, det_x, AY - AH, det_end - det_x, AH,
         RGBColor(0xFF, 0xF0, 0xE0))
add_textbox(s, det_x + Inches(0.05), AY - AH + Inches(0.05), Inches(1.7), Inches(0.3),
    "RetryGuard\ndetection window\n(≥ N × 30s above\nthreshold)", font_size=Pt(8), color=ACCENT)

# RetryGuard activation marker
act_x = det_x + (det_end - det_x) * 0.45
dline(s, act_x, AY - AH + Inches(0.05), act_x, AY, ACCENT, Pt(2))
add_textbox(s, act_x + Inches(0.05), cap_y - Inches(0.3), Inches(1.5), Inches(0.3),
    "Retries\nDISABLED", font_size=Pt(8.5), bold=True, color=ACCENT)

# Time axis labels
t_labels = ["0", "30s", "60s", "90s", "2m", "3m"]
for i, lbl in enumerate(t_labels):
    tx = AX + AW * i / (len(t_labels) - 1)
    dline(s, tx, AY, tx, AY + Inches(0.06), GRAY, Pt(1))
    add_textbox(s, tx - Inches(0.2), AY + Inches(0.05), Inches(0.4), Inches(0.22),
        lbl, font_size=Pt(7.5), color=GRAY, align=PP_ALIGN.CENTER)

note(s, "The core experiment. Only prolonged overload (ρ > 1 sustained for N intervals of ~30s) activates RetryGuard. Brief spikes do not.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Scenario: Targeted Bottleneck
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Targeted Bottleneck", "§7 — Load scenario 3 of 6")
add_section_label(s, "SECTION 7 OF 9 (3/6)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.3),
    "WHAT HAPPENS IN THE SYSTEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(1.0),
    "Full call chain exercises at normal RPS. One specific downstream service (e.g., Checkout or a mid-chain service) is constrained — reduced replica count or CPU limit — so it reaches ρ > 1 even under TopFull's throttled entry rate. TopFull detects and throttles APIs routing through it at entry. But after TopFull admits a request, the constrained service may still reject it, and its upstream caller retries via Istio — invisible to TopFull's proxy.",
    font_size=Pt(13), color=DARK, wrap=True)

divider_line(s, Inches(3.1))

add_textbox(s, Inches(0.4), Inches(3.25), Inches(12.5), Inches(0.3),
    "WHAT IT TESTS", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.55), Inches(12.5), Inches(0.75),
    "RetryGuard, reading Istio metrics directly at the bottleneck service, sees the rejection rate and suppresses per-service retries. Does this reduce load at the bottleneck faster and more directly than TopFull's top-down throttling alone? Does the benefit propagate upward through the call chain? Directly analogous to the RetryGuard Bookinfo case study (Sec. 6.2).",
    font_size=Pt(13), color=DARK, wrap=True)

divider_line(s, Inches(4.45))

add_textbox(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.3),
    "OPEN QUESTIONS ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.95), Inches(12.5), Inches(1.3),
    ["Topology beneficiaries — does RetryGuard's benefit concentrate at the constrained service?",
     "Chain propagation — do resource savings at the bottleneck propagate upward?",
     "Controller interaction — does per-service suppression complement TopFull's top-down throttling?"],
    font_size=Pt(13))

note(s, "The Bookinfo analogy in Sec. 6.2: Reviews service (slow HPA) vs. Product service (fast HPA). Checkout in Online Boutique is the analogous hub service.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Scenario: Topology Position Comparison
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Topology Position Comparison", "§7 — Load scenario 4 of 6")
add_section_label(s, "SECTION 7 OF 9 (4/6)")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.55),
    "Three separate Targeted Bottleneck runs — same load, same methodology — varying only where in the Online Boutique call chain the constrained service sits. Tests whether RetryGuard's suppression value scales with fan-out width and topology depth.",
    font_size=Pt(13), color=DARK, wrap=True)

positions = [
    ("(1) Gateway-adjacent / shallow sub-tree",
     "e.g., Recommendation or ProductCatalog. Called directly by Frontend; few or no downstream dependencies. TopFull's entry-level routing sees this bottleneck most directly. Retry amplification is limited to Frontend → service."),
    ("(2) Hub / sub-tree root",
     "e.g., Checkout. Fans out to Cart, Shipping, Currency, ProductCatalog, Email, and Payment. Creates bidirectional retry amplification: Frontend retries Checkout (upward) while Checkout retries all six downstream callers simultaneously (downward). Expected to show the most severe amplification."),
    ("(3) Deep leaf",
     "e.g., Email or Payment. No downstream dependencies; reachable only through Checkout. TopFull's top-down signal here is most attenuated — the bottleneck is invisible at the entry until Checkout itself degrades."),
]

y = Inches(2.35)
for title, body in positions:
    add_textbox(s, Inches(0.4), y, Inches(7.7), Inches(0.28),
                title, font_size=Pt(11), bold=True, color=ACCENT)
    add_textbox(s, Inches(0.4), y + Inches(0.3), Inches(7.7), Inches(0.7),
                body, font_size=Pt(12), color=DARK, wrap=True)
    y += Inches(1.1)

add_textbox(s, Inches(0.4), Inches(5.65), Inches(7.7), Inches(0.3),
    "OPEN QUESTIONS ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(5.95), Inches(7.7), Inches(0.45),
    "Topology position sensitivity, topology beneficiaries, chain propagation",
    font_size=Pt(13), color=DARK)

if os.path.exists(ARCH_IMG):
    s.shapes.add_picture(ARCH_IMG, Inches(8.3), Inches(1.4), Inches(4.6), Inches(5.5))
    add_textbox(s, Inches(8.3), Inches(6.95), Inches(4.6), Inches(0.35),
                "Annotate: circle ProductCatalog (1), Checkout (2), Email/Payment (3)",
                font_size=Pt(9), color=GRAY, italic=True)

note(s, "The hub position (Checkout) creates bidirectional amplification — suppressing retries at Checkout simultaneously relieves all six downstream callers.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Scenario: Re-enable Interval Tuning
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Re-enable Interval Tuning", "§7 — Load scenario 5 of 6")
add_section_label(s, "SECTION 7 OF 9 (5/6)")

LW = Inches(6.8)
add_textbox(s, Inches(0.4), Inches(1.65), LW, Inches(0.3),
    "WHAT HAPPENS IN THE SYSTEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), LW, Inches(0.65),
    "The Sustained Overload scenario is run multiple times, holding all parameters constant and varying only RetryGuard's re-enable interval (the Interval parameter in Algorithm 1). Test values: 10s, 20s, 30s (paper default, Sec. 6.2), 60s.",
    font_size=Pt(13), color=DARK, wrap=True)

add_textbox(s, Inches(0.4), Inches(2.75), LW, Inches(0.3),
    "WHAT IT TESTS", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(3.1), LW, Inches(2.3),
    ["Too short (10s): premature re-enabling before bottleneck clears, potentially re-triggering overload.",
     "Too long (60s): keeps retries suppressed after bottleneck clears, slowing throughput recovery.",
     "The paper's 30s default was validated without a co-running top-down controller. With TopFull's RL at 1s, recovery dynamics may be faster or more oscillatory. Does the optimal interval shift?"],
    font_size=Pt(12.5))

add_textbox(s, Inches(0.4), Inches(5.55), LW, Inches(0.3),
    "OPEN QUESTIONS ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(5.85), LW, Inches(0.4),
    "Interval parameter sensitivity, combined equilibrium",
    font_size=Pt(13), color=DARK)

vert_bar(s, Inches(7.35))

# ── Right: interval comparison bars ──────────────────────────────────────────
dpanel_bg(s, Inches(7.55), Inches(1.5), Inches(5.4), Inches(5.6))
panel_title(s, Inches(7.65), Inches(1.55), Inches(5.2),
    "  RE-ENABLE INTERVAL — TRADE-OFF")

intervals = [
    ("10 s",  0.72, ACCENT, "too short?\npremature"),
    ("20 s",  0.88, GREEN,  "candidate"),
    ("30 s",  1.0,  BLUE,   "paper\ndefault"),
    ("60 s",  0.78, GRAY,   "too long?\nslow recovery"),
]
bar_label_w = Inches(0.7)
bar_max_w   = Inches(3.5)
bar_h       = Inches(0.52)
bar_gap     = Inches(0.88)
bx          = Inches(8.2)
by          = Inches(2.25)
for label, rel, col, note_txt in intervals:
    bw = bar_max_w * rel
    add_textbox(s, bx - bar_label_w - Inches(0.05), by + Inches(0.1),
                bar_label_w, Inches(0.35), label, font_size=Pt(11), bold=True,
                color=DARK, align=PP_ALIGN.RIGHT)
    add_rect(s, bx, by, bw, bar_h, col)
    add_textbox(s, bx + bw + Inches(0.08), by + Inches(0.08),
                Inches(1.2), Inches(0.45),
                note_txt, font_size=Pt(9), color=col)
    by += bar_gap

add_textbox(s, Inches(7.65), Inches(6.1), Inches(5.2), Inches(0.55),
    "Bars = illustrative relative throughput recovery.\nActual values measured from experiment data.",
    font_size=Pt(9), color=GRAY, italic=True)

note(s, "Parameter sensitivity study. The hypothesis: TopFull's faster feedback loop changes the optimal Interval value.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Scenario: Attack Traffic
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Scenario: Attack Traffic", "§7 — Load scenario 6 of 6  ·  Extension, time-permitting")
add_section_label(s, "SECTION 7 OF 9 (6/6)")

add_rect(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.38), ACCENT)
add_textbox(s, Inches(0.5), Inches(1.68), Inches(12.3), Inches(0.32),
    "TIME-PERMITTING EXTENSION — only if Scenarios 1–5 are complete",
    font_size=Pt(12), bold=True, color=WHITE)

add_textbox(s, Inches(0.4), Inches(2.2), Inches(12.5), Inches(0.3),
    "WHAT HAPPENS IN THE SYSTEM", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(2.5), Inches(12.5), Inches(0.7),
    "A malicious burst-DDoS traffic pattern simulates an attacker attempting to exploit retry amplification. Short, concentrated bursts push one or more services into ρ > 1 briefly, relying on retry amplification to sustain the overload longer than the attack burst itself.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(3.35))

add_textbox(s, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.3),
    "WHAT IT TESTS", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.8), Inches(12.5), Inches(0.7),
    "Does hostile load trip RetryGuard at the wrong time — suppressing retries for legitimate requests when the high rejection rate is caused by an attacker rather than genuine overload? Or does RetryGuard correctly blunt retry storms caused by the attack without disrupting healthy service flows?",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(4.65))

add_textbox(s, Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.3),
    "OPEN QUESTION ADDRESSED", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.45),
    "Adversarial resilience — the RetryGuard paper (2025) notes DDoS amplification mitigation as a secondary benefit. This scenario tests that claim in the TopFull+Kubernetes context.",
    font_size=Pt(13), color=DARK, wrap=True)

note(s, "Source: PRESENTATION-GUIDE.md §7, RetryGuard.pdf (DDoS discussion). Mark clearly as time-permitting.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Metrics (Layers 1 & 2)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Metrics", "§8 — Three layers of measurement (1/2)")
add_section_label(s, "SECTION 8 OF 9 (1/2)")

# Layer 1
add_rect(s, Inches(0.4), Inches(1.65), Inches(0.08), Inches(1.5), BLUE)
add_textbox(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.3),
    "LAYER 1 — System & API Performance", font_size=Pt(12), bold=True, color=BLUE)
add_textbox(s, Inches(0.6), Inches(1.9), Inches(5.5), Inches(0.28),
    "Source: TopFull's metric_collector.py → CSVs in logs/", font_size=Pt(10.5), color=GRAY, italic=True)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.0),
    ["Goodput and latency per API (getcart, getproduct, postcheckout …) — primary outcome. Latency SLO: 1 second (as in TopFull paper).",
     "Rejection rate per API — the signal RetryGuard reads to decide whether to suppress retries.",
     "Retries per request — the most direct measure. Compare vs. RetryGuard paper Table 1 (Istio: 0.31 → 0.01)."],
    font_size=Pt(12.5))

divider_line(s, Inches(3.45))

# Layer 2
add_rect(s, Inches(0.4), Inches(3.6), Inches(0.08), Inches(1.55), ACCENT)
add_textbox(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(0.3),
    "LAYER 2 — Infrastructure Resource Usage", font_size=Pt(12), bold=True, color=ACCENT)
add_textbox(s, Inches(0.6), Inches(3.85), Inches(5.5), Inches(0.28),
    "Source: cAdvisor via resource_collector.py", font_size=Pt(10.5), color=GRAY, italic=True)
add_bullets(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(1.0),
    ["CPU consumption and memory limits per pod — tracks whether retry suppression actually frees resources at the service level.",
     "Pod instance counts over time (num_instances.csv) — shows autoscaling response and whether over-scaling is prevented (key finding in paper Sec. 6.2)."],
    font_size=Pt(12.5))

divider_line(s, Inches(5.35))

add_textbox(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.55),
    "All data synthesized into comparative time-series charts: baseline run vs. RetryGuard run, side by side across the same metrics.",
    font_size=Pt(13), color=DARK, italic=True)

note(s, "Layers 1 and 2 use TopFull's existing collection infrastructure — metric_collector.py, resource_collector.py.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Metrics (Layer 3)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Metrics (continued)", "§8 — Three layers of measurement (2/2)")
add_section_label(s, "SECTION 8 OF 9 (2/2)")

add_rect(s, Inches(0.4), Inches(1.65), Inches(0.08), Inches(2.1), NAVY)
add_textbox(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.3),
    "LAYER 3 — Controller Logic & State", font_size=Pt(12), bold=True, color=NAVY)
add_textbox(s, Inches(0.6), Inches(1.9), Inches(7), Inches(0.28),
    "Source: our RetryGuard script logs  +  TopFull's overload_detection.py",
    font_size=Pt(10.5), color=GRAY, italic=True)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.7),
    ["Which services had retries toggled off (and when) — ties controller decisions to topology beneficiaries question. Each toggle event logged with timestamp, service name, rejection rate that triggered it.",
     "Time-to-recovery: how long between retry suppression and re-enablement — shows the cool-down cycle and is directly relevant to interval parameter sensitivity.",
     "Business priority context from TopFull's overload_detection.py — which APIs were flagged as overloaded, so RetryGuard decisions can be cross-referenced with TopFull's simultaneous state."],
    font_size=Pt(12.5))

divider_line(s, Inches(4.1))

add_textbox(s, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.3),
    "CROSS-REFERENCING STRATEGY", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.65), Inches(12.5), Inches(1.6),
    ["RetryGuard toggle events (Layer 3) correlated with goodput changes (Layer 1).",
     "Resource drops (Layer 2) correlated with service-specific retry suppression (Layer 3).",
     "TopFull's overload_detection.py state correlated with RetryGuard's per-service decisions."],
    font_size=Pt(13))

add_textbox(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45),
    "Source: PRESENTATION-GUIDE.md §8, WORKPLAN.md (key metrics tables), PRESENTATION-ACTION-ITEMS.md §8",
    font_size=Pt(10), color=GRAY, italic=True)

note(s, "Layer 3 is the novel instrumentation we add. Layers 1 and 2 use existing TopFull infrastructure.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Timeline & Milestones (Gantt)
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Timeline & Milestones", "§9 — What happens and when")
add_section_label(s, "SECTION 9 OF 9")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.38),
    "A sequential view of the project phases. Dependencies are implied by the order.",
    font_size=Pt(13.5), color=DARK, italic=True)

# ── Gantt chart ───────────────────────────────────────────────────────────────
GX = Inches(3.1)    # Gantt area left edge
GY = Inches(2.2)    # Gantt area top
GW = Inches(9.8)    # Gantt area width
GH = Inches(4.4)    # Gantt area height
WEEKS = 4

# Week header columns
week_w = GW / WEEKS
for i in range(WEEKS):
    col = RGBColor(0xD8, 0xE8, 0xF8) if i % 2 == 0 else LTGRAY
    add_rect(s, GX + week_w * i, GY, week_w, Inches(0.42), col)
    add_textbox(s, GX + week_w * i, GY + Inches(0.07), week_w, Inches(0.3),
                f"Week {i+1}", font_size=Pt(11), bold=True, color=NAVY,
                align=PP_ALIGN.CENTER)

# Gantt rows
milestones = [
    ("Infrastructure setup\n(VMs, K8s, Istio, app running)",          0.0,  2.0, NAVY),
    ("Baseline experiment\n(TopFull running, default retries)",        1.5,  3.0, BLUE),
    ("RetryGuard implementation\n(Algorithm 1, Istio integration)",    2.0,  3.0, ACCENT),
    ("RetryGuard experiment\n(all load scenarios)",                    2.5,  4.0, BLUE),
    ("Evaluation & final report\n(comparison, time-series charts)",    3.0,  4.0, GREEN),
]

ROW_H = Inches(0.65)
ROW_GAP = Inches(0.1)
LBL_W = Inches(2.55)

for i, (lbl, start_w, end_w, col) in enumerate(milestones):
    ry = GY + Inches(0.42) + (ROW_H + ROW_GAP) * i

    # Row background
    add_rect(s, GX, ry, GW, ROW_H, LTGRAY)

    # Phase label (left of Gantt)
    add_textbox(s, Inches(0.3), ry + Inches(0.1), LBL_W - Inches(0.15), ROW_H - Inches(0.1),
                lbl, font_size=Pt(10), color=DARK, wrap=True)

    # Bar
    bar_x = GX + week_w * start_w
    bar_w = week_w * (end_w - start_w)
    bar = add_rect(s, bar_x, ry + Inches(0.1), bar_w, ROW_H - Inches(0.2), col)

    # Duration label inside bar
    dur_str = f"Week {start_w:.0f}–{end_w:.0f}" if end_w - start_w > 0.8 else ""
    if dur_str:
        add_textbox(s, bar_x + Inches(0.08), ry + Inches(0.2),
                    bar_w - Inches(0.1), Inches(0.3),
                    dur_str, font_size=Pt(9.5), bold=True, color=WHITE)

# Week divider lines
for i in range(1, WEEKS):
    dline(s, GX + week_w * i, GY, GX + week_w * i,
             GY + Inches(0.42) + (ROW_H + ROW_GAP) * len(milestones),
             MIDGRAY, Pt(1))

add_textbox(s, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
    "Source: WORKPLAN.md (all phases), PRESENTATION-ACTION-ITEMS.md §9",
    font_size=Pt(10), color=GRAY, italic=True)

note(s, "No phase numbers on slides — internal planning references. No blocker callouts — address them verbally if a mentor asks.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Summary & Deliverables
# ════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
set_bg(s, LTGRAY)
header_bar(s, "Summary & Deliverables", "What we are building, testing, and delivering")

add_textbox(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.3),
    "WHAT WE ARE BUILDING", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.55),
    "A full experimental stack: Kubernetes + Istio + TopFull overload controller + self-implemented RetryGuard, on Google Cloud Platform. Test application: Online Boutique (Google's 11-microservice demo app).",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(2.65))

add_textbox(s, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.3),
    "WHAT WE ARE TESTING", font_size=Pt(11), bold=True, color=BLUE)
add_textbox(s, Inches(0.4), Inches(3.1), Inches(12.5), Inches(0.55),
    "Whether RetryGuard (Algorithm 1, self-implemented from the RetryGuard paper) produces measurable improvement in goodput, latency, retry rate, and resource usage when layered on top of TopFull — across six load scenarios designed to answer eight open questions.",
    font_size=Pt(13.5), color=DARK, wrap=True)

divider_line(s, Inches(3.8))

add_textbox(s, Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.3),
    "DELIVERABLES", font_size=Pt(11), bold=True, color=BLUE)
add_bullets(s, Inches(0.4), Inches(4.3), Inches(12.5), Inches(1.7),
    ["Working Kubernetes + Istio + TopFull + RetryGuard experimental setup (reproducible, on GCP).",
     "Baseline vs. RetryGuard experiment data across all load scenarios — CSVs, logs, controller decision records.",
     "Evaluation report with time-series charts comparing goodput, latency, retries per request, resource usage, and autoscaler behavior — baseline vs. RetryGuard, side by side."],
    font_size=Pt(13.5))

divider_line(s, Inches(6.15))

add_textbox(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.4),
    "The outcome — whether RetryGuard helps significantly, helps only specific services, or barely moves the needle — is the point of the work.",
    font_size=Pt(13), color=NAVY, bold=True, wrap=True)

note(s, "Closing slide. The final sentence is from PRESENTATION-GUIDE.md — it captures the scientific stance.")


# ────────────────────────────────────────────────────────────────────────────
# Save
# ────────────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "RetryGuard-TopFull-ProjectPlan-v2.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
